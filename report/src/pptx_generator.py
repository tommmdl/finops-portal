"""
Gera o PowerPoint mensal de FinOps a partir do template.
"""
import io
import copy
from datetime import date
from pathlib import Path

MESES_PT = {
    "January": "Janeiro", "February": "Fevereiro", "March": "Março",
    "April": "Abril", "May": "Maio", "June": "Junho",
    "July": "Julho", "August": "Agosto", "September": "Setembro",
    "October": "Outubro", "November": "Novembro", "December": "Dezembro",
}

def _month_str_pt(d: date) -> str:
    en = d.strftime("%B %Y")
    month_en, year = en.split(" ")
    return f"{MESES_PT[month_en]} {year}"

import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ── Paleta e-core ─────────────────────────────────────────────────────────────
BLUE_DARK  = "#1F4E79"
BLUE_MID   = "#2E75B6"
BLUE_LIGHT = "#BDD7EE"
ORANGE     = "#F4B942"
GRAY       = "#595959"
WHITE      = "#FFFFFF"

CHART_COLORS = [BLUE_DARK, BLUE_MID, BLUE_LIGHT, ORANGE, "#70AD47",
                "#ED7D31", "#A5A5A5", "#4472C4", "#9E480E", "#636363"]

# Nomes curtos para serviços AWS nativos com nomes longos.
# Serviços não listados (ex: Marketplace) mantêm o nome original da API.
SERVICE_NAME_MAP = {
    "Amazon Simple Storage Service":                         "S3",
    "Amazon Elastic Compute Cloud - Compute":                "EC2-Instances",
    "EC2 - Other":                                           "EC2-Other",
    "Amazon Virtual Private Cloud":                          "VPC",
    "AmazonCloudWatch":                                      "CloudWatch",
    "Amazon Elastic File System":                            "Elastic File System",
    "Amazon Elastic Load Balancing":                         "Elastic Load Balancing",
    "Amazon FSx":                                            "FSx",
    "AWS Config":                                            "Config",
    "Amazon Relational Database Service":                    "RDS",
    "Amazon Elastic Container Service":                      "ECS",
    "Amazon Route 53":                                       "Route 53",
    "AWS Key Management Service":                            "KMS",
    "Amazon Redshift":                                       "Redshift",
    "Savings Plans for AWS Compute usage":                   "Savings Plans for Compute usage",
    "Amazon EC2 Container Registry (ECR)":                   "EC2 Container Registry (ECR)",
    "AWS Identity and Access Management Access Analyzer":    "IAM Access Analyzer",
    "Amazon Managed Workflows for Apache Airflow":           "MWAA",
    "Amazon Elastic MapReduce":                              "EMR",
    "AWS Lambda":                                            "Lambda",
    "Amazon DynamoDB":                                       "DynamoDB",
    "Amazon SageMaker":                                      "SageMaker",
    "Amazon Bedrock":                                        "Bedrock",
    "AWS Secrets Manager":                                   "Secrets Manager",
    "AWS CloudTrail":                                        "CloudTrail",
    "AWS Glue":                                              "Glue",
}


# ── Utilitários de gráfico ────────────────────────────────────────────────────

def _chart_to_image(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def _bar_chart_evolucao(totals: dict) -> io.BytesIO:
    """Gráfico de barras: evolução total de custos 6 meses."""
    months = list(totals.keys())
    values = list(totals.values())

    fig, ax = plt.subplots(figsize=(8, 4.5), facecolor="white")
    bars = ax.bar(months, values, color=BLUE_MID, width=0.55, zorder=3)

    ax.set_facecolor("white")
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(
        lambda x, _: f"${x:,.0f}"))
    ax.tick_params(axis="x", labelsize=9, rotation=15)
    ax.tick_params(axis="y", labelsize=9)
    ax.spines[["top", "right"]].set_visible(False)
    ax.yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)

    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + max(values) * 0.01,
                f"${val:,.0f}", ha="center", va="bottom", fontsize=8, fontweight="bold",
                color=BLUE_DARK)

    ax.set_title("Evolução de Custos AWS — 6 Meses", fontsize=11,
                 fontweight="bold", color=BLUE_DARK, pad=12)
    fig.tight_layout()
    return _chart_to_image(fig)


def _bar_chart_top10(top10: pd.DataFrame, last3_months: list[str]) -> io.BytesIO:
    """Eixo X = meses; dentro de cada mês, uma barra por serviço (cor única por serviço)."""
    services  = top10["Service"].tolist()
    n_svc     = len(services)
    n_months  = len(last3_months)
    if n_svc == 0 or n_months == 0:
        fig, ax = plt.subplots(figsize=(10, 5), facecolor="white")
        ax.text(0.5, 0.5, "Sem dados disponíveis", ha="center", va="center",
                transform=ax.transAxes, fontsize=14, color=GRAY)
        ax.axis("off")
        return _chart_to_image(fig)
    w         = 0.7 / n_svc
    x_base    = list(range(n_months))

    svc_colors = ["#1F4E79","#4472C4","#00B0F0","#375623","#7030A0",
                  "#FFC000","#ED7D31","#FF0000","#E6B8A2","#C00000"]

    fig, ax = plt.subplots(figsize=(10, 5), facecolor="white")
    ax.set_facecolor("white")

    for i, (svc, color) in enumerate(zip(services, svc_colors)):
        offsets = [xi + (i - n_svc / 2 + 0.5) * w for xi in x_base]
        vals    = [float(top10.loc[top10["Service"] == svc, m].values[0])
                   if m in top10.columns else 0 for m in last3_months]
        ax.bar(offsets, vals, width=w * 0.92, label=svc, color=color, zorder=3)

    ax.set_xticks(x_base)
    ax.set_xticklabels(last3_months, fontsize=10)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda v, _: f"${v:,.2f}"))
    ax.tick_params(axis="y", labelsize=9)
    ax.spines[["top", "right"]].set_visible(False)
    ax.yaxis.grid(True, linestyle="--", alpha=0.4, zorder=0)
    ax.legend(fontsize=7.5, loc="lower center",
              bbox_to_anchor=(0.5, -0.32), ncol=5, frameon=False)

    fig.tight_layout()
    return _chart_to_image(fig)


# ── Helpers de PPTX ───────────────────────────────────────────────────────────

def _replace_picture(slide, shape_idx: int, image_buf: io.BytesIO):
    """Substitui uma shape PICTURE pelo conteúdo de image_buf."""
    old = slide.shapes[shape_idx]
    left, top, width, height = old.left, old.top, old.width, old.height
    sp = old._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(image_buf, left, top, width, height)


def _set_text(shape, text: str):
    """Substitui o texto de um TextFrame preservando a formatação do primeiro run."""
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""
    if tf.paragraphs:
        para = tf.paragraphs[0]
        if para.runs:
            para.runs[0].text = text
        else:
            para.add_run().text = text


def _set_cell_text(cell, text: str, align=PP_ALIGN.CENTER):
    """
    Atualiza o texto de uma célula PRESERVANDO a formatação original
    (cor, negrito, tamanho). Não usa cell.text= para evitar reset de estilo.
    """
    tf = cell.text_frame
    # Mantém só o primeiro parágrafo
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[len(tf.paragraphs) - 1]._p
        p.getparent().remove(p)
    para = tf.paragraphs[0]
    para.alignment = align
    # Mantém só o primeiro run
    while len(para.runs) > 1:
        r = para.runs[len(para.runs) - 1]._r
        r.getparent().remove(r)
    if para.runs:
        para.runs[0].text = text
    else:
        para.add_run().text = text


def _update_table(table, data: list[list]):
    """Atualiza uma tabela existente redimensionando linhas se necessário,
    PRESERVANDO a formatação original de cada célula."""
    n_rows_needed = len(data)
    n_rows_have   = len(table.rows)

    # Clonar última linha de dados para adicionar se precisar de mais
    if n_rows_needed > n_rows_have:
        last_tr = table.rows[n_rows_have - 1]._tr
        for _ in range(n_rows_needed - n_rows_have):
            new_tr = copy.deepcopy(last_tr)
            last_tr.getparent().append(new_tr)

    # Remover linhas excedentes
    while len(table.rows) > n_rows_needed:
        tr = table.rows[len(table.rows) - 1]._tr
        tr.getparent().remove(tr)

    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            if c >= len(table.columns):
                break
            _set_cell_text(table.cell(r, c), str(val))


def _fmt_brl(value: float) -> str:
    s = f"{abs(value):,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    return f"${s}" if value >= 0 else f"-${s}"


def _fmt_reais(value: float) -> str:
    s = f"{abs(value):,.2f}".replace(",", "X").replace(".", ",").replace("X", ".")
    return f"R${s}" if value >= 0 else f"-R${s}"


def _rebuild_economia_table(slide, shape_idx: int,
                             utilization_df_ytd: pd.DataFrame,
                             usd_brl: float, reference: date):
    """Atualiza a tabela de Economia Acumulada (3 colunas x 4 linhas)."""
    shape = slide.shapes[shape_idx]
    table = shape.table

    sp_usd    = float(utilization_df_ytd["NetSavings"].sum()) \
                if not utilization_df_ytd.empty and "NetSavings" in utilization_df_ytd.columns \
                else 0.0
    rds_usd   = 0.0   # RI savings ainda não coletados
    total_usd = sp_usd + rds_usd

    data = [
        [f"Economia Acumulada {reference.year}", "Economia $",         "Economia R$"],
        ["Compute Savings Plans",                 _fmt_brl(sp_usd),     _fmt_reais(sp_usd  * usd_brl)],
        ["RDS",                                   _fmt_brl(rds_usd),    _fmt_reais(rds_usd * usd_brl)],
        ["Total",                                 _fmt_brl(total_usd),  _fmt_reais(total_usd * usd_brl)],
    ]
    _update_table(table, data)


def _rebuild_oscillation_table(slide, table_shape_idx: int,
                                osc_df: pd.DataFrame, month_cols: list[str]):
    """Reconstrói a tabela de oscilações com dados dinâmicos."""
    shape = slide.shapes[table_shape_idx]
    table = shape.table

    header = ["Service"] + month_cols + ["Média", "Aumento"]
    rows = [header]
    for _, row in osc_df.iterrows():
        vals = [float(row[m]) if m in row.index else 0.0 for m in month_cols]
        avg  = sum(vals) / len(vals) if vals else 0
        inc  = (vals[-1] / avg - 1) if avg != 0 else 0
        r    = [row["Service"]]
        r   += [_fmt_brl(v) for v in vals]
        r   += [_fmt_brl(avg), f"{inc:.2%}".replace(".", ",")]
        rows.append(r)

    _update_table(table, rows)


def _abbrev_month(label: str) -> str:
    """'Dezembro 2025' → 'Dez/25' — evita overflow nas células da tabela."""
    parts = label.split(" ")
    if len(parts) == 2:
        return f"{parts[0][:3]}/{parts[1][2:]}"
    return label


def _rebuild_coverage_table(slide, shape_idx: int, coverage_df: pd.DataFrame,
                             utilization_df: pd.DataFrame, ri_coverage_df: pd.DataFrame,
                             reference: date, n_months: int = 3):
    """Reconstrói a tabela de cobertura e utilização."""
    from dateutil.relativedelta import relativedelta

    # Meses fixos baseados no reference — independente do que tem dados no DynamoDB
    months = []
    d = date(reference.year, reference.month, 1)
    for _ in range(n_months):
        months.insert(0, _month_str_pt(d))
        d = d - relativedelta(months=1)

    shape = slide.shapes[shape_idx]
    table = shape.table
    months_abbrev = [_abbrev_month(m) for m in months]

    def _pct(df: pd.DataFrame, mes: str, col: str) -> str:
        if df.empty or col not in df.columns:
            return "0%"
        val = df[df["Mês"] == mes][col]
        # CE API já retorna em %, ex: 86.23 — não usar :.0% (multiplica por 100)
        return f"{val.values[0]:.0f}%" if len(val) else "0%"

    header = ["Cobertura"] + months_abbrev + ["", "Utilização"] + months_abbrev

    # Linha Compute Savings Plans — usa dados de SP
    sp_row  = ["Compute Savings Plans"]
    sp_row += [_pct(coverage_df,    m, "CoveragePercentage")    for m in months]
    sp_row += ["", "Compute Savings Plans"]
    sp_row += [_pct(utilization_df, m, "UtilizationPercentage") for m in months]

    # Linha RDS — usa dados de RI (riCoverage), independente do SP
    rds_row  = ["RDS"]
    rds_row += [_pct(ri_coverage_df, m, "CoveragePercentage")    for m in months]
    rds_row += ["", "RDS"]
    rds_row += [_pct(ri_coverage_df, m, "UtilizationPercentage") for m in months]

    print(f"[debug] tabela colunas: {len(table.columns)}")
    print(f"[debug] header length:  {len(header)}")
    print(f"[debug] header: {header}")
    _update_table(table, [header, sp_row, rds_row])


# ── Entry point ───────────────────────────────────────────────────────────────

def generate(
    costs_df: pd.DataFrame,
    coverage_df: pd.DataFrame,
    utilization_df: pd.DataFrame,
    ri_coverage_df: pd.DataFrame,
    utilization_df_ytd: pd.DataFrame,
    recommendations: list[dict],
    client_name: str,
    reference: date,
    usd_brl: float,
    template_path: Path,
    output_dir: Path,
) -> Path:
    month_str  = _month_str_pt(reference)
    month_cols = [c for c in costs_df.columns if c != "Service"]

    # Garante dtype numérico independente da origem dos dados (DynamoDB Decimal, etc.)
    for col in month_cols:
        costs_df[col] = pd.to_numeric(costs_df[col], errors="coerce").fillna(0.0)

    # Aplica nomes curtos — apenas serviços AWS nativos mapeados; demais mantêm nome original
    costs_df["Service"] = costs_df["Service"].map(
        lambda s: SERVICE_NAME_MAP.get(s, s)
    )

    print(f"[pptx_generator] costs_df shape: {costs_df.shape}, columns: {costs_df.columns.tolist()}")

    prs = Presentation(str(template_path))

    # ── Slide 1: Capa — só troca o nome do cliente, preserva toda formatação ──
    slide1 = prs.slides[0]
    for shape in slide1.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                # Troca qualquer nome de cliente conhecido pelo nome correto
                for known in ["Wilson Sons", "Delta Energia", "Easy Carros",
                               "Compliance", "Rediseg", "Ubots"]:
                    if known in run.text:
                        run.text = run.text.replace(known, client_name)

    # ── Slide 2: Evolução 6 meses (substituir gráfico) ────────────────────────
    slide2 = prs.slides[1]
    totals = {m: costs_df[m].sum() for m in month_cols}
    chart_buf = _bar_chart_evolucao(totals)

    # Encontra a picture maior (o gráfico) — maior área
    pictures = [(i, s) for i, s in enumerate(slide2.shapes)
                if str(s.shape_type) == "PICTURE (13)" and s.width > 2000000]
    if pictures:
        idx = max(pictures, key=lambda x: x[1].width * x[1].height)[0]
        _replace_picture(slide2, idx, chart_buf)

    # ── Slide 3: Top 10 serviços (substituir gráfico) ─────────────────────────
    slide3 = prs.slides[2]

    top10  = costs_df.nlargest(10, month_cols[-1]) if month_cols else pd.DataFrame()
    last3  = month_cols[-3:]
    print(f"[pptx_generator] top10 shape: {top10.shape}")
    if not top10.empty and last3:
        chart_buf2 = _bar_chart_top10(top10, last3)
        pictures3 = [(i, s) for i, s in enumerate(slide3.shapes)
                     if str(s.shape_type) == "PICTURE (13)" and s.width > 2000000]
        if pictures3:
            idx3   = max(pictures3, key=lambda x: x[1].width * x[1].height)[0]
            old3   = slide3.shapes[idx3]
            # Move o gráfico para baixo para não cobrir o subtítulo
            TOP_OFFSET = Inches(0.35)
            left3, top3, w3, h3 = old3.left, old3.top + TOP_OFFSET, old3.width, old3.height - TOP_OFFSET
            old3._element.getparent().remove(old3._element)
            slide3.shapes.add_picture(chart_buf2, left3, top3, w3, h3)
    else:
        print("[pptx_generator] WARNING: top10 vazio, slide 3 mantém imagem original")

    # ── Slide 4: Maior oscilação (atualizar tabela) ───────────────────────────
    # Apenas serviços com gasto real no último mês (>= $10) e com aumento efetivo
    slide4   = prs.slides[3]
    df_copy  = costs_df.copy()
    df_copy["_avg"] = df_copy[month_cols].mean(axis=1)
    df_copy["_inc"] = df_copy.apply(
        lambda r: (r[month_cols[-1]] / r["_avg"] - 1) if r["_avg"] != 0 else 0, axis=1
    )
    MIN_SPEND = 10.0
    osc_df = (
        df_copy[
            (df_copy[month_cols[-1]] >= MIN_SPEND) &   # elimina centavos
            (df_copy["_inc"] > 0)                       # só aumentos reais
        ]
        .nlargest(10, "_inc")[["Service"] + month_cols]
        .reset_index(drop=True)
    )

    table_shapes = [(i, s) for i, s in enumerate(slide4.shapes) if s.has_table]
    if table_shapes:
        tbl_idx = table_shapes[0][0]
        _rebuild_oscillation_table(slide4, tbl_idx, osc_df, month_cols)
        # Centraliza a tabela verticalmente abaixo do título (~1.2cm do topo)
        tbl_shape  = slide4.shapes[tbl_idx]
        title_btm  = Inches(1.2)
        avail      = prs.slide_height - title_btm
        tbl_shape.top = int(title_btm + (avail - tbl_shape.height) / 2)

    # ── Slide 5: Cobertura e Utilização (manual — não tocar se vazio) ────────
    if not coverage_df.empty:
        slide5 = prs.slides[4]
        cov_tables = [(i, s) for i, s in enumerate(slide5.shapes) if s.has_table]
        if cov_tables:
            _rebuild_coverage_table(slide5, cov_tables[0][0], coverage_df, utilization_df, ri_coverage_df, reference)

        # Tabela de Economia Acumulada YTD (segunda tabela do slide 5 — shape_idx=3)
        if len(cov_tables) >= 2:
            _rebuild_economia_table(slide5, cov_tables[1][0], utilization_df_ytd, usd_brl, reference)

    # ── Slide 6: Pontos de Atenção (atualizar texto) ──────────────────────────
    slide6 = prs.slides[5]
    if recommendations:
        pontos_text = "Itens que requerem atenção:\n"
        for rec in recommendations[:5]:
            pontos_text += f"\n• {rec['Recurso']} — Economia: ${rec['Economia']:,.2f}/mês"
        for shape in slide6.shapes:
            if shape.has_text_frame and "Itens que requerem" in shape.text_frame.text:
                _set_text(shape, pontos_text)

    # ── Salvar ────────────────────────────────────────────────────────────────
    filename = f"{client_name} _ Report Mensal {month_str}.pptx"
    path = output_dir / filename
    prs.save(str(path))
    return path
